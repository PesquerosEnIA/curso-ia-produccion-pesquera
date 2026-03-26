"""
Generador de presentaciones branded para el Curso IA Sector Pesquero
UTN FRCh · PesquerosEnIA

Genera .pptx para Clases 1, 4 y 10 con identidad visual del curso.
Uso: python3 scripts/gen_presentaciones.py
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

# ── Paleta de colores del curso ───────────────────────────────────────────────
UTNAZUL    = RGBColor(0,   70,  127)   # Azul institucional UTN
UTNCELESTE = RGBColor(0,  153,  204)   # Celeste UTN
ARIELBLUE  = RGBColor(44,  95,  124)   # Azul personal Ariel
ARIELGREEN = RGBColor(64, 145,  108)   # Verde Ariel
TECHGOLD   = RGBColor(212, 160,  23)   # Dorado Dr. González
PESCATEAL  = RGBColor(5,  150,  105)   # Teal Soraya
PESCAOCEAN = RGBColor(59, 130,  246)   # Ocean Damian
WHITE      = RGBColor(255, 255, 255)
DARKTEXT   = RGBColor(30,   30,   50)
MIDGRAY    = RGBColor(100, 110, 125)
LIGHTBG    = RGBColor(245, 247, 250)
ACCENTLINE = RGBColor(180, 195, 210)

# ── Dimensiones widescreen 16:9 ───────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

BASE = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
LOGO_UTN          = os.path.join(BASE, "propuesta", "logo_utn_frch_negro.png")
LOGO_UTN_40       = os.path.join(BASE, "propuesta", "logo_utn_frch_40aniv.png")
LOGO_PESQUEROS    = os.path.join(BASE, "propuesta", "logo_pesqueros_full.png")
LOGO_PESQ_SYMBOL  = os.path.join(BASE, "propuesta", "logo_pesqueros_symbol.png")
LOGO_ARIEL_FULL   = os.path.join(BASE, "propuesta", "logo_ariel_full.png")
LOGO_ARIEL_SYMBOL = os.path.join(BASE, "propuesta", "logo_ariel_symbol.png")


# ── Primitivas de dibujo ──────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill, alpha=None, line_color=None, line_width=None):
    from pptx.util import Pt
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if alpha is not None:
        shape.fill.fore_color.theme_color = None
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, x, y, w, h, size=18, color=WHITE, bold=False,
                align=PP_ALIGN.LEFT, font="Calibri", wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    return txBox


def add_bullet_textbox(slide, bullets, x, y, w, h, size=16, color=DARKTEXT,
                       accent=UTNCELESTE, font="Calibri", line_spacing=1.15):
    """Agrega un cuadro de texto con bullets (lista de strings o tuplas (texto, nivel))."""
    from pptx.oxml.ns import qn
    from pptx.util import Pt
    import copy

    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True

    first = True
    for item in bullets:
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()

        p.level = level
        # Espaciado
        p.space_before = Pt(4 if level == 0 else 2)
        p.space_after = Pt(2)

        run = p.add_run()
        run.text = ("   " * level) + ("• " if level == 0 else "  · ") + text
        run.font.size = Pt(size - level * 1.5)
        run.font.color.rgb = color
        run.font.name = font
        run.font.bold = (level == 0)

    return txBox


def add_logo(slide, logo_path, x, y, h):
    """Agrega logo desde PNG manteniendo proporción."""
    if not os.path.exists(logo_path):
        return None
    from PIL import Image as PILImage
    try:
        img = PILImage.open(logo_path)
        iw, ih = img.size
        ratio = iw / ih
        w = int(h * ratio)
        pic = slide.shapes.add_picture(logo_path, x, y, width=Emu(w), height=Emu(h))
        return pic
    except Exception:
        # sin PIL: ancho fijo
        pic = slide.shapes.add_picture(logo_path, x, y, height=Emu(h))
        return pic


def add_logo_placeholder(slide, label, x, y, w, h, color=UTNCELESTE):
    """Placeholder de color para logos en PDF (no disponibles en PNG)."""
    rect = add_rect(slide, x, y, w, h, color)
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(9)
    run.font.color.rgb = WHITE
    run.font.bold = True
    run.font.name = "Calibri"
    return rect


def divider_line(slide, x, y, w, color=UTNCELESTE, thickness=Pt(1.5)):
    from pptx.util import Pt
    line = slide.shapes.add_connector(1, x, y, x + w, y)  # STRAIGHT connector
    line.line.color.rgb = color
    line.line.width = thickness
    return line


# ── Layouts de slides ─────────────────────────────────────────────────────────

def make_cover_slide(prs, title, subtitle, docentes, clase_num, accent=UTNAZUL):
    """Slide de portada con fondo azul, logos reales y datos del docente."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # ── Fondo ────────────────────────────────────────────────────────────────
    add_rect(slide, 0, 0, W, H, accent)

    # Franja diagonal decorativa (rectángulo rotado simulado con dos rectángulos)
    # Banda de acento en diagonal — simulada con un bloque semitransparente
    dark_accent = RGBColor(
        max(0, accent[0] - 30),
        max(0, accent[1] - 10),
        max(0, accent[2] - 20)
    )
    add_rect(slide, W - Inches(4.5), 0, Inches(4.5), H, dark_accent)

    # Franja dorada superior
    add_rect(slide, 0, 0, W, Inches(0.10), TECHGOLD)
    # Franja dorada inferior
    add_rect(slide, 0, H - Inches(0.10), W, Inches(0.10), TECHGOLD)

    # ── Franja inferior para datos del curso ─────────────────────────────────
    add_rect(slide, 0, H - Inches(1.15), W, Inches(1.05),
             RGBColor(0, 0, 0))  # negro semitransparente
    # Superponer un rect con el color de acento a menor opacidad
    strip = add_rect(slide, 0, H - Inches(1.15), W, Inches(1.05),
                     RGBColor(max(0, accent[0]-15), max(0, accent[1]+20), min(255, accent[2]+40)))

    # ── Logos en la franja superior ──────────────────────────────────────────
    logo_y  = Inches(0.17)
    logo_h  = Inches(0.80)

    # Logo UTN FRCh (izquierda)
    if os.path.exists(LOGO_UTN):
        add_logo(slide, LOGO_UTN, Inches(0.40), logo_y, Emu(int(logo_h)))

    # Logo PesquerosEnIA (derecha) — ahora con PNG real
    if os.path.exists(LOGO_PESQUEROS):
        pesq_h = Inches(0.65)
        pic = add_logo(slide, LOGO_PESQUEROS, W - Inches(3.6), logo_y + Inches(0.08), Emu(int(pesq_h)))
    else:
        add_logo_placeholder(slide, "PesquerosEnIA",
                             W - Inches(2.2), logo_y, Inches(2.0), logo_h, PESCATEAL)

    # Logo Ariel (símbolo pequeño, junto al de Pesqueros)
    if os.path.exists(LOGO_ARIEL_SYMBOL):
        ariel_h = Inches(0.65)
        add_logo(slide, LOGO_ARIEL_SYMBOL, W - Inches(1.2), logo_y + Inches(0.08), Emu(int(ariel_h)))

    # ── Separador bajo logos ──────────────────────────────────────────────────
    add_rect(slide, Inches(0.40), Inches(1.18), W - Inches(0.80), Inches(0.035), UTNCELESTE)

    # ── Número de clase ───────────────────────────────────────────────────────
    add_textbox(slide, f"CLASE {clase_num}",
                Inches(0.45), Inches(1.35), Inches(4), Inches(0.42),
                size=12, color=UTNCELESTE, bold=True, font="Calibri")

    # ── Título principal ──────────────────────────────────────────────────────
    add_textbox(slide, title,
                Inches(0.45), Inches(1.75), W - Inches(5.2), Inches(2.5),
                size=34, color=WHITE, bold=True, font="Calibri", align=PP_ALIGN.LEFT)

    # ── Subtítulo ─────────────────────────────────────────────────────────────
    add_textbox(slide, subtitle,
                Inches(0.45), Inches(4.15), W - Inches(5.2), Inches(0.75),
                size=17, color=UTNCELESTE, bold=False, font="Calibri", align=PP_ALIGN.LEFT)

    # ── Decoración lateral derecha (área oscura) ──────────────────────────────
    add_textbox(slide, "IA · Pesca · Datos",
                W - Inches(4.2), Inches(2.5), Inches(3.8), Inches(1.0),
                size=22, color=RGBColor(180, 210, 240), bold=False,
                font="Calibri", align=PP_ALIGN.CENTER, italic=True)

    add_textbox(slide, "UTN FRCh + PesquerosEnIA",
                W - Inches(4.2), Inches(3.4), Inches(3.8), Inches(0.6),
                size=13, color=RGBColor(150, 190, 230), bold=False,
                font="Calibri", align=PP_ALIGN.CENTER)

    # ── Datos en franja inferior ──────────────────────────────────────────────
    add_textbox(slide, f"Docentes: {docentes}",
                Inches(0.45), H - Inches(1.05), Inches(9), Inches(0.48),
                size=14, color=WHITE, bold=True, font="Calibri")

    add_textbox(slide,
                "Inteligencia Artificial Aplicada a la Produccion Pesquera  |  UTN FRCh  |  2026",
                Inches(0.45), H - Inches(0.60), W - Inches(0.9), Inches(0.45),
                size=10, color=RGBColor(190, 220, 250), bold=False, font="Calibri",
                align=PP_ALIGN.LEFT)

    return slide


def make_agenda_slide(prs, items, accent=UTNAZUL):
    """Slide de agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Fondo suave
    add_rect(slide, 0, 0, W, H, LIGHTBG)
    # Barra de color izquierda
    add_rect(slide, 0, 0, Inches(0.22), H, accent)
    # Header
    add_rect(slide, 0, 0, W, Inches(1.4), accent)
    add_textbox(slide, "AGENDA DE LA CLASE", Inches(0.4), Inches(0.08), W - Inches(1), Inches(0.55),
                size=13, color=UTNCELESTE, bold=True, font="Calibri")
    add_textbox(slide, "Lo que vamos a ver hoy", Inches(0.4), Inches(0.55), W - Inches(1), Inches(0.75),
                size=26, color=WHITE, bold=True, font="Calibri")

    # Items de agenda
    content_y = Inches(1.65)
    col_w = (W - Inches(1.2)) / 2
    for i, item in enumerate(items):
        col = i % 2
        row = i // 2
        ix = Inches(0.5) + col * (col_w + Inches(0.2))
        iy = content_y + row * Inches(0.72)
        # Número
        add_rect(slide, ix, iy + Inches(0.05), Inches(0.45), Inches(0.45), accent)
        add_textbox(slide, str(i + 1), ix, iy + Inches(0.02), Inches(0.45), Inches(0.45),
                    size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font="Calibri")
        # Texto
        add_textbox(slide, item, ix + Inches(0.55), iy, col_w - Inches(0.65), Inches(0.6),
                    size=15, color=DARKTEXT, bold=False, font="Calibri")

    _add_footer(slide, accent)
    return slide


def make_content_slide(prs, title, subtitle_note, bullets, accent=UTNAZUL, icon=""):
    """Slide de contenido con header de color y bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, LIGHTBG)
    add_rect(slide, 0, 0, Inches(0.22), H, accent)
    # Header
    add_rect(slide, 0, 0, W, Inches(1.5), accent)
    if subtitle_note:
        add_textbox(slide, subtitle_note.upper(), Inches(0.4), Inches(0.08), W - Inches(1), Inches(0.45),
                    size=11, color=UTNCELESTE, bold=True, font="Calibri")
    add_textbox(slide, title, Inches(0.4), Inches(0.4) if subtitle_note else Inches(0.35),
                W - Inches(1.2), Inches(1.0), size=26, color=WHITE, bold=True, font="Calibri")

    # Bullets
    bullet_items = []
    for b in bullets:
        if isinstance(b, str):
            bullet_items.append((b, 0))
        else:
            bullet_items.append(b)

    add_bullet_textbox(slide, bullet_items,
                       Inches(0.5), Inches(1.65),
                       W - Inches(0.8), H - Inches(2.0),
                       size=17, color=DARKTEXT, accent=accent)

    _add_footer(slide, accent)
    return slide


def make_table_slide(prs, title, subtitle_note, headers, rows, accent=UTNAZUL):
    """Slide con tabla de datos."""
    from pptx.util import Inches, Pt
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, LIGHTBG)
    add_rect(slide, 0, 0, Inches(0.22), H, accent)
    add_rect(slide, 0, 0, W, Inches(1.5), accent)
    if subtitle_note:
        add_textbox(slide, subtitle_note.upper(), Inches(0.4), Inches(0.08), W - Inches(1), Inches(0.45),
                    size=11, color=UTNCELESTE, bold=True, font="Calibri")
    add_textbox(slide, title, Inches(0.4), Inches(0.4) if subtitle_note else Inches(0.35),
                W - Inches(1.2), Inches(1.0), size=26, color=WHITE, bold=True, font="Calibri")

    # Tabla
    n_cols = len(headers)
    n_rows = len(rows) + 1  # +1 para header
    tbl_x = Inches(0.5)
    tbl_y = Inches(1.65)
    tbl_w = W - Inches(0.8)
    tbl_h = H - Inches(2.2)

    col_w = tbl_w / n_cols
    row_h = tbl_h / n_rows

    # Header de tabla
    add_rect(slide, tbl_x, tbl_y, tbl_w, row_h, accent)
    for c, hdr in enumerate(headers):
        add_textbox(slide, hdr, tbl_x + c * col_w + Inches(0.08), tbl_y + Inches(0.04),
                    col_w - Inches(0.1), row_h - Inches(0.08),
                    size=13, color=WHITE, bold=True, font="Calibri", wrap=True)

    # Filas de datos
    for r, row in enumerate(rows):
        bg = LIGHTBG if r % 2 == 0 else WHITE
        add_rect(slide, tbl_x, tbl_y + (r + 1) * row_h, tbl_w, row_h, bg)
        for c, cell in enumerate(row):
            is_first = (c == 0)
            add_textbox(slide, str(cell),
                        tbl_x + c * col_w + Inches(0.08),
                        tbl_y + (r + 1) * row_h + Inches(0.04),
                        col_w - Inches(0.1), row_h - Inches(0.08),
                        size=13, color=accent if is_first else DARKTEXT,
                        bold=is_first, font="Calibri", wrap=True)

    # Línea de borde inferior
    add_rect(slide, tbl_x, tbl_y + n_rows * row_h - Inches(0.02), tbl_w, Inches(0.04), accent)

    _add_footer(slide, accent)
    return slide


def make_quote_slide(prs, quote, author="", accent=UTNAZUL):
    """Slide con cita destacada."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, accent)
    add_rect(slide, 0, 0, Inches(0.5), H, TECHGOLD)
    add_rect(slide, 0, H - Inches(0.12), W, Inches(0.12), TECHGOLD)

    # Comilla decorativa
    add_textbox(slide, "\u201c", Inches(0.8), Inches(0.5), Inches(2), Inches(2),
                size=120, color=RGBColor(30, 90, 150), bold=True, font="Georgia")

    # Cita
    add_textbox(slide, quote, Inches(0.8), Inches(1.8), W - Inches(1.4), Inches(3.5),
                size=26, color=WHITE, bold=False, font="Georgia", italic=True,
                align=PP_ALIGN.LEFT)

    if author:
        add_textbox(slide, f"— {author}", Inches(0.8), H - Inches(1.8), W - Inches(1.4), Inches(0.6),
                    size=16, color=UTNCELESTE, bold=True, font="Calibri")

    return slide


def make_two_col_slide(prs, title, subtitle_note, left_content, right_content,
                       accent=UTNAZUL, left_title="", right_title=""):
    """Slide en dos columnas."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, LIGHTBG)
    add_rect(slide, 0, 0, Inches(0.22), H, accent)
    add_rect(slide, 0, 0, W, Inches(1.5), accent)
    if subtitle_note:
        add_textbox(slide, subtitle_note.upper(), Inches(0.4), Inches(0.08), W - Inches(1), Inches(0.45),
                    size=11, color=UTNCELESTE, bold=True, font="Calibri")
    add_textbox(slide, title, Inches(0.4), Inches(0.4) if subtitle_note else Inches(0.35),
                W - Inches(1.2), Inches(1.0), size=26, color=WHITE, bold=True, font="Calibri")

    col_w = (W - Inches(1.1)) / 2 - Inches(0.1)
    col_h = H - Inches(2.1)

    # Columna izquierda
    add_rect(slide, Inches(0.4), Inches(1.55), col_w, col_h, WHITE)
    if left_title:
        add_rect(slide, Inches(0.4), Inches(1.55), col_w, Inches(0.45), accent)
        add_textbox(slide, left_title, Inches(0.5), Inches(1.58), col_w - Inches(0.1), Inches(0.4),
                    size=12, color=WHITE, bold=True, font="Calibri")
        ty = Inches(2.1)
    else:
        ty = Inches(1.65)
    add_bullet_textbox(slide, [(b, 0) for b in left_content],
                       Inches(0.5), ty, col_w - Inches(0.1), col_h - Inches(0.55),
                       size=15, color=DARKTEXT)

    # Separador
    add_rect(slide, Inches(0.4) + col_w + Inches(0.05), Inches(1.55), Inches(0.04), col_h, ACCENTLINE)

    # Columna derecha
    rx = Inches(0.4) + col_w + Inches(0.2)
    add_rect(slide, rx, Inches(1.55), col_w, col_h, WHITE)
    if right_title:
        add_rect(slide, rx, Inches(1.55), col_w, Inches(0.45), UTNCELESTE)
        add_textbox(slide, right_title, rx + Inches(0.1), Inches(1.58), col_w - Inches(0.1), Inches(0.4),
                    size=12, color=WHITE, bold=True, font="Calibri")
        ty2 = Inches(2.1)
    else:
        ty2 = Inches(1.65)
    add_bullet_textbox(slide, [(b, 0) for b in right_content],
                       rx + Inches(0.1), ty2, col_w - Inches(0.1), col_h - Inches(0.55),
                       size=15, color=DARKTEXT)

    _add_footer(slide, accent)
    return slide


def make_closing_slide(prs, title, subtitle, repo, clase_siguiente="", accent=UTNAZUL):
    """Slide de cierre con logos reales."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, accent)

    # Panel oscuro decorativo derecho
    dark_accent = RGBColor(max(0,accent[0]-30), max(0,accent[1]-10), max(0,accent[2]-20))
    add_rect(slide, W - Inches(3.8), 0, Inches(3.8), H, dark_accent)

    add_rect(slide, 0, 0, W, Inches(0.10), TECHGOLD)
    add_rect(slide, 0, H - Inches(0.10), W, Inches(0.10), TECHGOLD)

    # Logos alineados arriba
    logo_h = Inches(0.80)
    if os.path.exists(LOGO_UTN):
        add_logo(slide, LOGO_UTN, Inches(0.40), Inches(0.17), Emu(int(logo_h)))

    if os.path.exists(LOGO_PESQUEROS):
        add_logo(slide, LOGO_PESQUEROS, Inches(3.2), Inches(0.22), Emu(int(Inches(0.65))))

    if os.path.exists(LOGO_ARIEL_SYMBOL):
        add_logo(slide, LOGO_ARIEL_SYMBOL, Inches(5.6), Inches(0.22), Emu(int(Inches(0.65))))

    add_rect(slide, Inches(0.40), Inches(1.18), W - Inches(4.2), Inches(0.035), UTNCELESTE)

    add_textbox(slide, title, Inches(0.40), Inches(1.40), W - Inches(4.5), Inches(1.8),
                size=32, color=WHITE, bold=True, font="Calibri", align=PP_ALIGN.LEFT)

    add_textbox(slide, subtitle, Inches(0.40), Inches(3.15), W - Inches(4.5), Inches(0.75),
                size=16, color=UTNCELESTE, bold=False, font="Calibri", align=PP_ALIGN.LEFT)

    add_rect(slide, Inches(0.40), Inches(4.05), Inches(5.5), Inches(0.035), TECHGOLD)

    add_textbox(slide, repo, Inches(0.40), Inches(4.15), Inches(8.5), Inches(0.55),
                size=14, color=RGBColor(190, 225, 255), bold=False, font="Consolas",
                align=PP_ALIGN.LEFT)

    if clase_siguiente:
        add_rect(slide, Inches(0.40), Inches(4.85), Inches(8.5), Inches(0.90),
                 RGBColor(max(0,accent[0]-20), min(255,accent[1]+30), min(255,accent[2]+50)))
        add_textbox(slide, clase_siguiente, Inches(0.55), Inches(4.90), Inches(8.2), Inches(0.80),
                    size=13, color=WHITE, bold=False, font="Calibri", align=PP_ALIGN.LEFT)

    # Panel lateral derecho
    add_textbox(slide, "github.com/\nPesquerosEnIA",
                W - Inches(3.5), Inches(2.2), Inches(3.2), Inches(1.4),
                size=16, color=RGBColor(180, 215, 255), bold=False,
                font="Consolas", align=PP_ALIGN.CENTER)

    add_textbox(slide, "UTN FRCh  |  PesquerosEnIA  |  2026",
                Inches(0.40), H - Inches(0.52), W - Inches(4.5), Inches(0.42),
                size=11, color=RGBColor(160, 200, 240), bold=False, font="Calibri",
                align=PP_ALIGN.LEFT)

    return slide


def _add_footer(slide, accent=UTNAZUL):
    """Footer con nombre del curso, logo símbolo y línea de color."""
    add_rect(slide, 0, H - Inches(0.48), W, Inches(0.48), accent)
    # Logo símbolo PesquerosEnIA en el footer (izquierda)
    if os.path.exists(LOGO_PESQ_SYMBOL):
        add_logo(slide, LOGO_PESQ_SYMBOL,
                 Inches(0.15), H - Inches(0.44), Emu(int(Inches(0.38))))
    add_textbox(slide,
                "IA Aplicada a la Produccion Pesquera  |  UTN FRCh  |  PesquerosEnIA  |  2026",
                Inches(0.65), H - Inches(0.44), W - Inches(0.8), Inches(0.38),
                size=10, color=UTNCELESTE, bold=False, font="Calibri", align=PP_ALIGN.CENTER)


# ── Función principal ─────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def save(prs, path):
    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)
    print(f"  >> Guardado: {path}")


# ═══════════════════════════════════════════════════════════════════════════════
# CLASE 1 — Estrategia y Transformación Digital
# ═══════════════════════════════════════════════════════════════════════════════

def gen_clase01():
    print("\nGenerando Clase 1...")
    prs = new_prs()
    AC = UTNAZUL

    # Slide 1 — Portada
    make_cover_slide(prs,
        title="Estrategia y Transformación Digital en el Sector Pesquero",
        subtitle="Clase 1 — Inteligencia Artificial Aplicada a la Producción Pesquera",
        docentes="Ariel Giamportone · Soraya Corvalán",
        clase_num=1, accent=AC)

    # Slide 2 — Agenda
    make_agenda_slide(prs, [
        "El sector pesquero argentino en números",
        "La Cuarta Revolución Industrial llega al mar",
        "¿Qué es una Smart Fishery?",
        "5 casos de uso de IA concretos",
        "La brecha tecnológica y oportunidades",
        "PesquerosEnIA: de dónde venimos",
        "Actividad práctica: mapeá tu oportunidad",
    ], accent=AC)

    # Slide 3 — El mar argentino en números
    make_content_slide(prs,
        "Argentina: una potencia pesquera con enorme potencial digital",
        "El sector pesquero argentino en números",
        [
            "ZEE: ~1.000.000 km² de Mar Argentino",
            "Exportaciones: ~USD 2.100 millones/año",
            "Capturas totales: ~850.000 toneladas/año",
            "Empleo directo: ~60.000 personas (embarcados + tierra)",
            "~600 buques habilitados — puertos en Patagonia, Bs.As., TDF",
            "",
            "\"Con estos números, una mejora del 5% en eficiencia vía IA vale más de USD 100 millones.\"",
        ], accent=AC)

    # Slide 4 — Principales especies (tabla)
    make_table_slide(prs,
        "Las especies que mueven la economía pesquera argentina",
        "Principales especies y su valor",
        ["Especie", "Participación", "Destino principal"],
        [
            ["Merluza hubbsi", "~40% del volumen", "UE, Brasil"],
            ["Calamar illex", "~30% del volumen", "Asia, UE"],
            ["Langostino patagónico", "~15% del valor", "UE, EE.UU."],
            ["Vieira patagónica", "En crecimiento", "Asia"],
        ], accent=AC)

    # Slide 5 — Industria 4.0
    make_content_slide(prs,
        "Industria 4.0: ya está pasando, también en el mar",
        "La Cuarta Revolución Industrial",
        [
            "IoT: sensores que capturan datos 24/7 (temperatura, posición, captura)",
            "Big Data: millones de registros de AIS, partes de pesca, sensores oceanográficos",
            "Inteligencia Artificial: algoritmos que aprenden y predicen",
            "Conectividad: Starlink ya está en barcos de altura",
            "",
            "¿Qué lo hace diferente? Las máquinas ya no solo ejecutan — también aprenden y deciden.",
            "",
            "Analogía: la agricultura de precisión ya usa drones, imágenes satelitales y ML para soja.",
            "El sector pesquero puede hacer exactamente lo mismo.",
        ], accent=AC)

    # Slide 6 — Particularidades del sector (tabla)
    make_table_slide(prs,
        "El mar tiene sus propias reglas",
        "¿Qué hace diferente al sector pesquero?",
        ["Característica", "Desafío", "Oportunidad IA"],
        [
            ["Recurso invisible y variable", "No podés ver los peces", "Predicción por datos ambientales"],
            ["Ambiente hostil", "Sensores deben ser robustos", "Monitoreo remoto satelital"],
            ["Cadena fría crítica", "Trazabilidad compleja", "Sensores + blockchain"],
            ["Regulación estricta", "Carga de reportes", "Automatización de partes de pesca"],
        ], accent=AC)

    # Slide 7 — Smart Fishery
    make_content_slide(prs,
        "Smart Fishery = Datos + Conectividad + Algoritmos + Decisión",
        "El concepto de Smart Fishery",
        [
            "Sensores/AIS/Satélites → [DATOS] → Algoritmos IA → [DECISIÓN] → Capitán / Planta / Regulador",
            "",
            "Pilar 1: Eficiencia operativa — menos combustible, más captura por marea",
            "Pilar 2: Sostenibilidad — pesca dentro de los límites del stock",
            "Pilar 3: Trazabilidad — del anzuelo al plato, verificable",
            "Pilar 4: Cumplimiento — reportes automáticos, menos carga burocrática",
        ], accent=AC)

    # Slide 8 — Casos globales líderes
    make_content_slide(prs,
        "El mundo ya avanzó — esto es lo que hacen los líderes",
        "Casos globales líderes",
        [
            "Noruega: flota con transmisión en tiempo real + IA para optimización de redes (Kongsberg)",
            "Islandia: 90% de la flota reporta captura digitalmente en tiempo real al gobierno",
            "Chile: imágenes SAR para monitoreo de acuicultura y detección de flota irregular",
            "Global Fishing Watch: monitorea 60.000 barcos con ML — detección de pesca ilegal >95% precisión",
            "",
            "Estos países tienen flota comparable a Argentina en tamaño.",
            "No son 'países del futuro' — ya lo implementaron.",
        ], accent=AC)

    # Slides 9–13 — Casos de uso
    make_content_slide(prs,
        "IA que predice dónde están los peces",
        "Caso de uso 1: Predicción de zonas de pesca",
        [
            "Problema: el capitán decide a qué zona ir con información limitada",
            "Solución: modelos de ML que procesan SST, clorofila, corrientes e historial de capturas",
            "",
            "Impacto demostrado:",
            ("Reducción 20-30% en días de búsqueda improductiva", 1),
            ("Ahorro de combustible: ~15-20% por marea", 1),
            ("Aumento del rendimiento (ton/día en mar)", 1),
            "",
            "Datos: Copernicus Marine Service, NOAA, historial de partes de pesca",
            "En la Clase 6 construiremos este modelo con Random Forest.",
        ], accent=AC)

    make_content_slide(prs,
        "Cada barco tiene una huella digital en el mar",
        "Caso de uso 2: Monitoreo satelital y AIS",
        [
            "AIS (Automatic Identification System): obligatorio >300 GT",
            "Transmite: posición, velocidad, rumbo, ID del barco cada 2-10 segundos",
            "Global Fishing Watch: 50 TB de datos AIS procesados por día",
            "",
            "VMS (Vessel Monitoring System): obligatorio pesqueros argentinos",
            "Gestiona Subsecretaría de Pesca · Transmisión satelital cada 30-60 minutos",
            "",
            "Dark vessels: barcos que apagan el AIS → detectados por imágenes SAR + IA",
        ], accent=AC)

    make_content_slide(prs,
        "Cámaras que ven mejor y más rápido que el ojo humano",
        "Caso de uso 3: Visión artificial en planta",
        [
            "Clasificación por especie a >100 peces/minuto (merluza / polaca / castañeta)",
            "Medición automática de talla individual — reemplaza medición manual",
            "Detección de parásitos o defectos en filetes",
            "Evaluación de frescura por color y textura",
            "",
            "Tecnología disponible hoy: Marel, Baader, TriVision",
            "Ya operan en plantas de procesamiento en Argentina",
        ], accent=AC)

    make_content_slide(prs,
        "Del anzuelo al plato — verificable con un escaneo",
        "Caso de uso 4: Trazabilidad digital",
        [
            "El recorrido: Barco → Puerto → Frigorífico → Planta → Exportador → Supermercado",
            "",
            "Cada paso registrado con timestamp y firma digital",
            "Código QR en el producto final",
            "El consumidor escanea y ve: barco, zona, fecha, certificaciones",
            "",
            "En Argentina: el Sistema de Trazabilidad Pesquera (SiTP) en desarrollo",
            "En Europa: requisito obligatorio para importar desde 2025+",
            "",
            "No es solo un requisito regulatorio — es una ventaja comercial.",
        ], accent=AC)

    make_content_slide(prs,
        "Cada litro de gasoil cuenta — la IA lo sabe",
        "Caso de uso 5: Optimización de flota",
        [
            "Rutas optimizadas: algoritmos con corrientes, vientos, zonas → ahorro 10-20% combustible",
            "Velocidad óptima: la ley cúbica del combustible naval — cada nudo extra cuesta el triple",
            "Mantenimiento predictivo: sensores en motores detectan anomalías antes de la avería",
            "Logística portuaria: coordinación de desembarco, frigorífico y planta",
            "",
            "Ejemplo: flota de 20 arrastreros de altura en Argentina",
            "Ahorro proyectado: USD 800.000–1.200.000/año en combustible",
            "",
            "En la Clase 8 vamos a construir este modelo juntos.",
        ], accent=AC)

    # Slide 14 — La brecha tecnológica
    make_content_slide(prs,
        "¿Por qué el sector pesquero argentino todavía no despegó?",
        "La brecha tecnológica: diagnóstico honesto",
        [
            "Formación: pocos profesionales del sector con competencias en datos e IA",
            "Conectividad: internet satelital recién llega a costos accesibles",
            "Inversión: costo percibido alto para empresas medianas",
            "Idioma: casi todo el material de calidad está en inglés",
            "Cultura: 'siempre lo hicimos así' — resistencia al cambio",
            "Regulación: el marco regulatorio no siempre acompaña la innovación",
            "",
            "Estas barreras existen en todas las industrias primarias.",
            "La buena noticia: todas son superables.",
        ], accent=AC)

    # Slide 15 — Las señales de cambio
    make_content_slide(prs,
        "El viento está cambiando — y hay que estar preparado",
        "Las señales de cambio",
        [
            "Pampa Azul: programa MINCYT con foco en economía del conocimiento marina",
            "Starlink en barcos: ya en operación en flotas argentinas de altura",
            "INIDEP abre datos: datasets históricos de gran valor, disponibles para análisis",
            "Empresas líderes adoptando trazabilidad para cumplir requisitos de exportación UE",
            "PesquerosEnIA: comunidad activa de profesionales del sector aprendiendo IA",
            "",
            "No se trata de si el sector se va a digitalizar.",
            "Se trata de si vamos a ser protagonistas o espectadores de esa transformación.",
        ], accent=AC)

    # Slide 16 — PesquerosEnIA
    make_content_slide(prs,
        "PesquerosEnIA — IA para el sector pesquero, desde el sector pesquero",
        "Nuestra comunidad",
        [
            "2019: I CONIPE — minicurso Industria 4.0 (8 horas) + mesa redonda",
            "2025: CONIPE 2025 — curso alfabetización en IA (8 módulos, Puerto Madryn)",
            "2026: Este curso — 23 horas de formación aplicada",
            "",
            "Principios:",
            ("Contenido en español, contextualizado al sector", 1),
            ("Acceso abierto (GPL-3.0)", 1),
            ("Colaboración entre pares", 1),
            ("Desde el sector, para el sector", 1),
            "",
            "github.com/PesquerosEnIA",
        ], accent=AC)

    # Slide 17 — El equipo docente
    make_content_slide(prs,
        "Quiénes están detrás del curso",
        "El equipo docente",
        [
            "Ariel Giamportone — Ing. Pesquero UTN FRCh | Data Science & IA | MBA Log. y Op. | Intervalor Data, Madrid",
            "Soraya Corvalán — Ing. Pesquera | INIDEP | FAO | Profesora Asociada UTN FRCh | 1ª graduada Ing. Pesquera ARG",
            "Damian Adolfo Giacone — Lic. Sistemas | Máster Industria 4.0 IEBS | Docente UGR | Ex-Alpesca S.A.",
            "Dr. Juan D. González — Dr. Mat. Aplicada UBA | SENASA IA | CONICET | paquetes RMBC + ktaucenters en CRAN",
        ], accent=AC)

    # Slide 18 — Estructura del curso (tabla)
    make_table_slide(prs,
        "23 horas de IA aplicada al sector pesquero",
        "Estructura del curso",
        ["Clase", "Tema", "Docente"],
        [
            ["1 ← HOY", "Estrategia y Transformación Digital", "Ariel · Soraya"],
            ["2", "Fundamentos de IA y LLMs", "Damian"],
            ["3", "Prompting y Asistentes de IA", "Damian"],
            ["4", "Datos y Sensores del Dominio Pesquero", "Ariel · Soraya"],
            ["5", "Arquitectura de Datos y Big Data", "Damian"],
            ["6", "Machine Learning y Modelos Predictivos", "Ariel · Damian"],
            ["7", "Visión Artificial en Plantas", "Damian"],
            ["8", "Optimización y Agentes de IA", "Ariel · Damian"],
            ["9", "Visualización y Dashboards", "Damian"],
            ["10", "Cierre y Hoja de Ruta", "Ariel · Soraya"],
        ], accent=AC)

    # Slide 19 — Competencias
    make_content_slide(prs,
        "El horizonte: competencias concretas para el sector",
        "¿Qué van a poder hacer al terminar el curso?",
        [
            "Usar LLMs y prompting para automatizar reportes y análisis",
            "Acceder y explorar datos oceanográficos y de captura reales",
            "Construir modelos de ML para predecir zonas de pesca",
            "Entender y comunicar los resultados de un modelo de visión artificial",
            "Optimizar rutas y operaciones con herramientas de IA",
            "Armar dashboards de toma de decisiones para tu organización",
            "Diseñar un caso de uso de IA para tu contexto específico",
        ], accent=AC)

    # Slide 20 — Actividad práctica (tabla)
    make_table_slide(prs,
        "Actividad — 20 minutos",
        "Mapeá tu oportunidad de IA",
        ["Campo", "Tu respuesta"],
        [
            ["Mi rol en el sector", ""],
            ["El problema que identifico", ""],
            ["Datos existentes sobre este problema", ""],
            ["Tipo de IA (predicción / clasificación / optimización)", ""],
            ["Beneficio esperado", ""],
        ], accent=AC)

    # Slide 21 — Cierre y síntesis
    make_content_slide(prs,
        "Lo que vimos hoy",
        "Síntesis de la clase",
        [
            "El sector pesquero argentino es estratégico — enorme potencial de digitalización",
            "La IA no es magia — son algoritmos que aprenden de datos que ya existen en el sector",
            "Smart Fisheries integra sensores + conectividad + datos + algoritmos",
            "5 aplicaciones concretas: predicción de zonas, AIS, visión artificial, trazabilidad, optimización",
            "La brecha existe pero es superable — este curso es un paso concreto",
            "PesquerosEnIA: una comunidad que ya está construyendo este puente",
        ], accent=AC)

    # Slide 22 — Próxima clase
    make_closing_slide(prs,
        title="¡Nos vemos en la Clase 2!",
        subtitle="Fundamentos de IA y Modelos de Lenguaje (LLMs) · Docente: Damian Adolfo Giacone",
        repo="github.com/PesquerosEnIA/curso-ia-produccion-pesquera",
        clase_siguiente="Para la próxima: probá ChatGPT o Claude con una pregunta de tu sector pesquero.",
        accent=AC)

    out = os.path.join(BASE, "clase_01_estrategia_transformacion", "slides", "slides_clase01.pptx")
    save(prs, out)


# ═══════════════════════════════════════════════════════════════════════════════
# CLASE 4 — Datos y Sensores del Dominio Pesquero
# ═══════════════════════════════════════════════════════════════════════════════

def gen_clase04():
    print("\nGenerando Clase 4...")
    prs = new_prs()
    AC = ARIELBLUE

    make_cover_slide(prs,
        title="Datos y Sensores del Dominio Pesquero",
        subtitle="Clase 4 — De la señal del sensor al insight para la pesca",
        docentes="Ariel Giamportone · Soraya Corvalán",
        clase_num=4, accent=AC)

    make_agenda_slide(prs, [
        "Datos que genera un barco pesquero",
        "Variables oceanográficas clave: SST, clorofila, corrientes",
        "AIS y VMS: el rastro digital de la flota",
        "Global Fishing Watch",
        "Plataformas de datos abiertos",
        "Formatos: CSV, NetCDF, GeoJSON",
        "Notebook práctico: exploración de datos de la PCA",
    ], accent=AC)

    # Slide 3 — Datos del barco (tabla)
    make_table_slide(prs,
        "Un arrastrero genera datos 24/7 — ¿cuántos explotamos?",
        "¿Qué datos genera un barco pesquero en un día?",
        ["Sensor / Sistema", "Dato generado", "Frecuencia"],
        [
            ["GPS", "Posición (lat/lon)", "Continua"],
            ["AIS", "Posición + velocidad + identidad", "Cada 2-10 s"],
            ["VMS", "Posición satelital", "Cada 30 min"],
            ["Ecosonda", "Profundidad + cardúmenes", "Continua"],
            ["Sensor SST", "Temperatura del agua", "Cada pocos min"],
            ["Balanza", "Peso captura por lance", "Por lance"],
            ["Parte de pesca", "Especie, kg, zona, hora", "Por lance"],
        ], accent=AC)

    # Slide 4 — Datos externos
    make_content_slide(prs,
        "El contexto ambiental: datos que no vienen del barco pero son críticos",
        "Los datos externos que también importan",
        [
            "Imágenes satelitales: SST, clorofila-a, nivel del mar (cada 1-3 días, gratis)",
            "Modelos oceanográficos: corrientes, temperatura subsuperficial, salinidad",
            "Meteorología: viento, altura de ola, visibilidad (seguridad + eficiencia)",
            "Regulatorios: cuotas por especie, áreas de veda, límites de ZEE",
            "Biológicos (INIDEP): índices de abundancia, campañas de evaluación",
            "",
            "La predicción de zonas de pesca (Clase 6) integra exactamente estos datos",
            "con el historial de capturas del barco.",
        ], accent=AC)

    # Slide 5 — SST
    make_content_slide(prs,
        "La SST — el termómetro que guía la pesca",
        "Temperatura Superficial del Mar (SST)",
        [
            "Cada especie tiene su rango térmico preferencial:",
            ("Merluza hubbsi: 4–12°C", 1),
            ("Calamar illex: 8–16°C", 1),
            ("Langostino patagónico: 6–14°C", 1),
            "",
            "Los frentes oceánicos (mezcla Malvinas-Brasil) = alta productividad",
            "Resolución disponible: hasta ~1 km por día (satélites MODIS, Sentinel)",
            "",
            "Fuentes gratuitas: Copernicus Marine Service · NOAA CoastWatch · NASA Earthdata",
        ], accent=AC)

    # Slide 6 — Clorofila
    make_content_slide(prs,
        "Verde en el mar = alimento para los peces",
        "Clorofila-a: productividad del mar",
        [
            "Clorofila-a = fitoplancton = base de la cadena trófica marina",
            "Alta clorofila → zooplancton → anchoíta/juveniles → merluza/calamar",
            "Floración de primavera (sep-nov): explosión de productividad en golfos patagónicos",
            "Rango en la PCA: 0.1 – 30 mg/m³",
            "",
            "Combinación poderosa:",
            ("SST en rango óptimo + alta clorofila = zona de alta probabilidad de captura", 1),
        ], accent=AC)

    # Slide 7 — Corrientes
    make_content_slide(prs,
        "Malvinas vs. Brasil — el gran motor de la productividad pesquera argentina",
        "Corrientes oceánicas de la Plataforma Continental Argentina",
        [
            "Corriente de Malvinas: fría (~6°C), rica en nutrientes, fluye al norte por el talud",
            "Corriente de Brasil: cálida (~22°C), pobre en nutrientes, fluye al sur",
            "",
            "Frente de confluencia: ~38-40°S — zona de máxima diversidad y captura",
            "El frente varía estacionalmente → cambia la distribución de recursos",
            "",
            "Este es el 'corazón oceanográfico' de la pesquería argentina.",
            "Los modelos de ML van a aprender estos patrones.",
        ], accent=AC)

    # Slide 8 — AIS
    make_content_slide(prs,
        "AIS — cada barco tiene una huella en el mar",
        "AIS: el pasaporte digital del barco",
        [
            "MMSI: identificador único de 9 dígitos (Argentina: comienza con 701)",
            "Ejemplo de mensaje AIS:",
            ("MMSI: 701234567  |  Nombre: BP PATAGÓNICO", 1),
            ("Lat: -43.582°  |  Lon: -60.213°", 1),
            ("SOG: 3.2 kn  ← velocidad lenta → probablemente pescando", 1),
            ("Timestamp: 2026-01-15 03:42:17 UTC", 1),
            "",
            "SOG < 4 kn → probable arrastre (pescando)",
            "SOG > 7 kn → navegación de ida/vuelta a puerto",
            "Global Fishing Watch usa esta lógica (+ML) para clasificar 70.000 barcos",
        ], accent=AC)

    # Slide 9 — AIS vs VMS (tabla)
    make_table_slide(prs,
        "VMS — el sistema de control obligatorio para la flota argentina",
        "AIS vs VMS: comparación",
        ["Aspecto", "AIS", "VMS"],
        [
            ["Obligatoriedad", "Buques ≥300 GT", "Buques pesqueros ≥28 m"],
            ["Datos", "Públicos", "Reservados (regulatorio)"],
            ["Frecuencia", "2-10 segundos", "30-60 minutos"],
            ["Propósito", "Seguridad náutica", "Control pesquero"],
            ["Gestión Argentina", "Prefectura / ITU", "Subsecretaría de Pesca"],
        ], accent=AC)

    # Slide 10 — Global Fishing Watch
    make_content_slide(prs,
        "Ver el mar como nunca antes",
        "Global Fishing Watch",
        [
            "Monitorea ~70.000 barcos en tiempo real",
            "50 TB de datos AIS procesados por día con modelos de ML",
            "Clasifica: pesca · navegación · transbordo · buques oscuros",
            "Precisión >95% en clasificación de actividad pesquera",
            "Datos descargables para investigación (registro gratuito)",
            "Detecta pesca ilegal en zonas protegidas y vedas",
            "",
            "Para Argentina: permite ver actividad de flota extranjera en límites de ZEE",
            "y el esfuerzo pesquero histórico por zona.",
        ], accent=AC)

    # Slide 11 — Plataformas de datos (tabla)
    make_table_slide(prs,
        "El sector tiene más datos gratuitos de los que imaginamos",
        "Plataformas de datos abiertos",
        ["Plataforma", "Qué tiene", "Acceso"],
        [
            ["Copernicus Marine", "SST, clorofila, corrientes, nivel del mar", "Registro gratuito"],
            ["NOAA ERDDAP", "SST, anomalías, temp. subsuperficial", "Sin registro"],
            ["INIDEP", "Capturas ARG, evaluaciones de stock", "Web pública"],
            ["Global Fishing Watch", "AIS, esfuerzo pesquero, rutas", "Registro gratuito"],
            ["NASA Earthdata", "MODIS, VIIRS, datos de satélite", "Registro gratuito"],
            ["GEBCO", "Batimetría global", "Sin registro"],
        ], accent=AC)

    # Slide 12 — Formatos de datos (tabla)
    make_table_slide(prs,
        "CSV, NetCDF, GeoJSON — el lenguaje de los datos oceánicos",
        "Formatos de datos: ¿qué vamos a encontrar?",
        ["Formato", "Descripción", "Cuándo aparece"],
        [
            ["CSV / Excel", "Tabular, filas y columnas", "Partes de pesca, producción planta"],
            ["NetCDF (.nc)", "Cubo 3D: lat × lon × tiempo", "SST, clorofila, corrientes satelitales"],
            ["GeoTIFF (.tif)", "Imagen georreferenciada", "SAR, imágenes satelitales"],
            ["Shapefile (.shp)", "Vectores geográficos", "Áreas de veda, límites ZEE"],
            ["JSON / GeoJSON", "Texto estructurado", "APIs de GFW, NOAA; datos AIS"],
        ], accent=AC)

    # Slide 13 — Transición al notebook
    make_content_slide(prs,
        "Ahora vamos a hacerlo nosotros",
        "Notebook práctico",
        [
            "Notebook: clase04_exploracion_datos_oceanograficos.ipynb",
            "Ejecutar en: Google Colab (sin instalación requerida)",
            "Link: github.com/PesquerosEnIA/curso-ia-produccion-pesquera",
            "",
            "Lo que vamos a hacer:",
            ("1. Dataset de SST realista de la PCA (lat -55° a -34°, lon -65° a -44°)", 1),
            ("2. Visualizar SST media anual y variación estacional", 1),
            ("3. Analizar clorofila-a y frentes productivos", 1),
            ("4. Explorar datos AIS simulados de un barco pesquero", 1),
            ("5. Cruzar datos de captura con variables ambientales", 1),
        ], accent=AC)

    # Slide 14 — SST PCA (visual placeholder)
    make_content_slide(prs,
        "La Plataforma Continental Argentina desde el satélite",
        "Mapa de SST — resultado del notebook",
        [
            "[Mapa generado por el notebook — SST media anual, lat -55° a -34°]",
            "",
            "Gradiente latitudinal claro: norte ~18°C → sur ~5°C",
            "Zona de frente (mezcla Malvinas-Brasil) visible en ~40°S",
            "Baja SST en invierno austral (jun-ago) → calamar migra al norte",
            "Alta SST en verano → langostino se activa en los golfos",
        ], accent=AC)

    # Slide 15 — AIS velocidad
    make_content_slide(prs,
        "¿Cómo saber si un barco está pescando sin verlo?",
        "Velocidad AIS como proxy de actividad",
        [
            "[Gráfico del notebook — trayectoria del barco coloreada por velocidad]",
            "",
            "Regla empírica para arrastrero:",
            ("Velocidad < 4 kn → probable arrastre (pescando)", 1),
            ("Velocidad 4-7 kn → transitando entre lances", 1),
            ("Velocidad > 7 kn → navegación de ida/vuelta a puerto", 1),
            "",
            "Este criterio es la base del algoritmo de detección de pesca de Global Fishing Watch.",
        ], accent=AC)

    # Slide 16 — Captura vs SST
    make_content_slide(prs,
        "¿A qué temperatura pesca más la merluza?",
        "Captura vs SST: el patrón que buscamos",
        [
            "[Scatter plot del notebook — SST vs toneladas de captura]",
            "",
            "El scatter muestra una relación no lineal",
            "Pico de captura en el rango 8–12°C → consistente con biología de merluza hubbsi",
            "Por debajo de 5°C o encima de 14°C → capturas significativamente menores",
            "",
            "Esta relación es la que el modelo de ML de la Clase 6 va a aprender.",
            "Los datos explorados hoy serán los predictores del modelo de mañana.",
        ], accent=AC)

    # Slide 17 — Flujo de trabajo (tabla)
    make_content_slide(prs,
        "De los datos crudos al insight — el workflow completo",
        "Síntesis: flujo de trabajo con datos pesqueros",
        [
            "Fuente de datos   →  Acceso          →  Procesamiento  →  Análisis",
            "─────────────────────────────────────────────────────────────────",
            "Copernicus (SST)  →  API Python      →  xarray/pandas  →  EDA + viz",
            "NOAA (batimetría) →  ERDDAP          →  pandas         →  correlación",
            "Partes de pesca   →  CSV             →  pandas         →  patrones",
            "AIS / VMS         →  GFW API         →  geopandas      →  actividad",
            "",
            "Todos los pasos son realizables en Python + Jupyter.",
            "El notebook completo de hoy está disponible en el repositorio.",
        ], accent=AC)

    # Slide 18 — Cierre
    make_closing_slide(prs,
        title="Lo que aprendimos · Lo que viene",
        subtitle="Clase 5: Arquitectura de Datos y Big Data · Docente: Damian Adolfo Giacone",
        repo="github.com/PesquerosEnIA/curso-ia-produccion-pesquera",
        clase_siguiente="Para la próxima: explorá globalfishingwatch.org/map — seleccioná el Mar Argentino.",
        accent=AC)

    out = os.path.join(BASE, "clase_04_datos_sensores", "slides", "slides_clase04.pptx")
    save(prs, out)


# ═══════════════════════════════════════════════════════════════════════════════
# CLASE 10 — Cierre, Conclusiones y Hoja de Ruta
# ═══════════════════════════════════════════════════════════════════════════════

def gen_clase10():
    print("\nGenerando Clase 10...")
    prs = new_prs()
    AC = PESCATEAL  # Color Soraya como protagonista del cierre

    make_cover_slide(prs,
        title="De los datos a las decisiones en el sector pesquero",
        subtitle="Clase 10 — Cierre, Conclusiones y Hoja de Ruta",
        docentes="Soraya Corvalán · Ariel Giamportone",
        clase_num=10, accent=AC)

    # Slide 2 — El camino recorrido
    make_content_slide(prs,
        "23 horas — un recorrido completo",
        "El camino recorrido",
        [
            "C1: Estrategia → C2-3: LLMs → C4-5: Datos → C6: ML → C7: Visión → C8: Optimiz. → C9: Dashboards → C10: Cierre",
            "+ 2 Charlas Dr. González: Acústica pesquera | Clustering avanzado",
            "",
            "Lo que construimos juntos:",
            ("Un mapa de la transformación digital del sector pesquero", 1),
            ("Competencias concretas en IA aplicada a la industria", 1),
            ("Una red de profesionales del sector con interés en IA", 1),
            "",
            "¿Qué fue lo que más les impactó del curso?",
        ], accent=AC)

    # Slide 3 — Cadena de valor del dato
    make_content_slide(prs,
        "Del GPS a la decisión estratégica — el poder de los datos",
        "La cadena de valor del dato pesquero",
        [
            "Sensor/GPS/Satélite → Python/Pandas → Modelo ML → Dashboard → Decisión",
            "(datos crudos)         (limpieza)      (predicción)  (visualiz.)  (capitán/gerente)",
            "",
            "Ejemplo concreto integrador:",
            ("Entrada: posición GPS + SST de Copernicus + historial de capturas", 1),
            ("Proceso: modelo Random Forest (construido en Clase 6)", 1),
            ("Salida: 'Zona B frente a Camarones — 78% de prob. de captura exitosa mañana'", 1),
            "",
            "Este es el hilo conductor de todo el curso.",
        ], accent=AC)

    # Slide 4 — Antes y después (tabla)
    make_table_slide(prs,
        "Antes y después — el impacto concreto de la IA",
        "¿Qué cambia con la IA?",
        ["Tarea", "Sin IA", "Con IA"],
        [
            ["Elegir zona de pesca", "Experiencia + radio + intuición", "Modelo predictivo con datos ambientales"],
            ["Redactar parte de pesca", "Manual, 30-45 minutos", "Asistente IA, 5 minutos (revisión humana)"],
            ["Clasificar peces en planta", "Operarios manuales, 100/min", "Visión artificial, 300/min consistente"],
            ["Anomalía en motor", "Avería → parada no planificada", "Alerta predictiva 2-3 semanas antes"],
            ["Reportar a la autoridad", "Formularios manuales", "Automatización con VMS/IA"],
        ], accent=AC)

    # Slide 5 — Cuadrante de priorización
    make_two_col_slide(prs,
        "Fruta baja: alta prioridad, bajo costo de implementación",
        "Lo que se puede implementar hoy",
        left_content=[
            "Predicción zonas pesca ★",
            "Dashboards de gestión ★",
            "LLMs para reportes ★",
            "Alertas de mantenimiento ★",
        ],
        right_content=[
            "Visión artificial en planta ●",
            "Trazabilidad blockchain ●",
            "Big Data en nube ●",
        ],
        left_title="Bajo costo / Alto impacto — IMPLEMENTAR AHORA",
        right_title="Alto costo / Alto impacto — MEDIANO PLAZO",
        accent=AC)

    # Slide 6 — Proyectos integradores
    make_content_slide(prs,
        "Sus proyectos — la IA en su propio contexto",
        "Proyectos integradores",
        [
            "Espacio para presentaciones del grupo (30-40 minutos)",
            "",
            "Estructura de cada presentación (5-7 minutos):",
            ("1. El problema identificado en su sector", 1),
            ("2. Los datos disponibles (o a conseguir)", 1),
            ("3. La solución IA propuesta", 1),
            ("4. El impacto esperado (cuantificado)", 1),
            ("5. Las barreras a superar y cómo enfrentarlas", 1),
        ], accent=AC)

    # Slide 7 — Hoja de ruta técnico-operativo
    make_two_col_slide(prs,
        "Si sos capitán, técnico o supervisor de planta...",
        "Hoja de ruta: perfil técnico-operativo",
        left_content=[
            "Próximos 30 días:",
            "Registrate en Claude/ChatGPT → probalo con tareas de tu trabajo",
            "Abrí el GitHub de PesquerosEnIA → explorá los notebooks",
            "Identificá UN dataset que tenés y que podría analizarse mejor",
        ],
        right_content=[
            "3-6 meses:",
            "Completá el curso Python básico de Kaggle Learn (gratis)",
            "Construí tu primer análisis de datos en Colab",
            "Armá un dashboard simple en Power BI con datos de producción",
            "Largo plazo: proponé un proyecto piloto de IA en tu organización",
        ],
        accent=AC)

    # Slide 8 — Hoja de ruta científico-académico
    make_two_col_slide(prs,
        "Si sos investigador, docente o profesional técnico...",
        "Hoja de ruta: perfil científico-académico",
        left_content=[
            "Próximos 30 días:",
            "Explorá datasets de Copernicus Marine Service",
            "Reproducí el notebook de Clase 6 con datos reales de INIDEP",
            "Publicá un análisis propio en GitHub",
        ],
        right_content=[
            "3-6 meses:",
            "Aplicá ML a tu investigación: predicción, clasificación o clustering",
            "Contribuí al repo de PesquerosEnIA con tu especialidad",
            "Considerá integrar IA en las materias que dictás",
            "Largo plazo: artículo académico + co-autoría en PesquerosEnIA",
        ],
        accent=AC)

    # Slide 9 — Hoja de ruta gestor/regulador
    make_two_col_slide(prs,
        "Si sos tomador de decisiones en empresa u organismo...",
        "Hoja de ruta: perfil gestor/regulador",
        left_content=[
            "Próximos 30 días:",
            "Identificá el dato más crítico que hoy no explotás bien",
            "Explorá globalfishingwatch.org/map — tu zona de interés",
            "Compartí los materiales del curso con tu equipo",
        ],
        right_content=[
            "3-6 meses:",
            "Levantá un proyecto piloto pequeño (dashboard, modelo básico)",
            "Evaluá soluciones de trazabilidad digital para exportación",
            "Formá un equipo interno con las competencias del curso",
            "Largo plazo: estrategia de datos para tu organización",
        ],
        accent=AC)

    # Slide 10 — PesquerosEnIA
    make_content_slide(prs,
        "Una comunidad que crece con cada uno de ustedes",
        "PesquerosEnIA: lo que construimos juntos",
        [
            "Lo que existe hoy:",
            ("3+ cursos impartidos: CONIPE 2019, CONIPE 2025, este", 1),
            ("Repositorios con notebooks de ML/DL pesquero en acceso abierto", 1),
            ("Red de profesionales del sector con competencias en IA", 1),
            "",
            "Lo que queremos construir:",
            ("Base de datos curada de datasets pesqueros argentinos abiertos", 1),
            ("Más notebooks y guías en español contextualizados al sector", 1),
            ("Casos de uso documentados de implementación de IA en empresas", 1),
            "",
            "Cómo sumarse: Star en GitHub · Fork + Pull Request · Compartir con colegas",
            "github.com/PesquerosEnIA",
        ], accent=AC)

    # Slide 11 — IA ética
    make_quote_slide(prs,
        quote="La IA no nos dice qué queremos lograr — eso lo decidimos nosotros como sociedad pesquera.",
        author="Reflexión de cierre del curso",
        accent=UTNAZUL)

    # Slide 12 — Recursos para seguir (tabla)
    make_table_slide(prs,
        "El ecosistema de aprendizaje que tienen disponible",
        "Recursos para seguir",
        ["Recurso", "Para qué", "URL"],
        [
            ["PesquerosEnIA GitHub", "Materiales del curso", "github.com/PesquerosEnIA"],
            ["Global Fishing Watch", "Datos y herramientas AIS", "globalfishingwatch.org"],
            ["Copernicus Marine", "Datos oceanográficos SST, clorofila", "marine.copernicus.eu"],
            ["Kaggle Learn", "Cursos gratuitos ML y Python", "kaggle.com/learn"],
            ["INIDEP", "Datos del Atlántico Sur", "inidep.edu.ar"],
            ["Google Colab", "Notebooks sin instalación", "colab.research.google.com"],
        ], accent=AC)

    # Slide 13 — El equipo docente
    make_content_slide(prs,
        "Lo que cada docente pone en este curso",
        "El equipo docente",
        [
            "Ariel Giamportone: 10 años de modelado predictivo aplicado al sector pesquero y acuícola.",
            "  Fundador de PesquerosEnIA. De la Patagonia a Intervalor Data, Madrid.",
            "",
            "Soraya Corvalán: gestión de recursos pesqueros en Argentina, INIDEP, FAO, Pampa Azul.",
            "  Primera graduada Ing. Pesquera Argentina. Profesora Asociada Concursada UTN FRCh.",
            "",
            "Damian Giacone: 20 años de tecnología aplicada al sector productivo.",
            "  Raíces en la industria pesquera patagónica (Alpesca S.A.), Máster IEBS, UGR.",
            "",
            "Dr. Juan D. González: matemática de frontera para problemas pesqueros reales.",
            "  SENASA IA, CONICET/DIIV-ARA, paquetes RMBC + ktaucenters en CRAN.",
        ], accent=AC)

    # Slide 14 — Certificación
    make_content_slide(prs,
        "Su certificado UTN FRCh",
        "Certificación del curso",
        [
            "Requisitos para certificación:",
            ("Asistencia: ≥75% de las clases", 1),
            ("Proyecto integrador: presentado en esta clase", 1),
            "",
            "Emite: Universidad Tecnológica Nacional, Facultad Regional Chubut",
            "Avala: PesquerosEnIA",
            "Tipo: Certificado de Extensión Universitaria",
            "",
            "Proceso administrativo: se informará por correo electrónico.",
        ], accent=AC)

    # Slide 15 — Mensaje de cierre (quote)
    make_quote_slide(prs,
        quote="La IA no va a reemplazar a los ingenieros pesqueros, biólogos, capitanes ni operadores de planta. Va a reemplazar a los que no saben trabajar con IA.",
        author="",
        accent=UTNAZUL)

    # Slide 16 — Cierre formal
    make_closing_slide(prs,
        title="¡Gracias por ser parte de PesquerosEnIA!",
        subtitle="El aprendizaje no termina aquí — comienza la aplicación.",
        repo="github.com/PesquerosEnIA/curso-ia-produccion-pesquera",
        clase_siguiente="Todos los materiales disponibles bajo licencia GPL-3.0",
        accent=AC)

    out = os.path.join(BASE, "clase_10_cierre", "slides", "slides_clase10.pptx")
    save(prs, out)


# ── Main ──────────────────────────────────────────────────────────────────────
if __name__ == "__main__":
    print("=" * 60)
    print("Generador de Presentaciones - Curso IA Sector Pesquero")
    print("=" * 60)
    gen_clase01()
    gen_clase04()
    gen_clase10()
    print("\nListo! 3 presentaciones generadas.")
    print(f"  clase_01.../slides/slides_clase01.pptx")
    print(f"  clase_04.../slides/slides_clase04.pptx")
    print(f"  clase_10.../slides/slides_clase10.pptx")
